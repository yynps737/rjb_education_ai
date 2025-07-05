"""
智能文档处理器，支持多种格式的文档解析和处理
"""
import os
import re
import json
import hashlib
from typing import List, Dict, Optional, Any, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
import logging
import mimetypes
from io import BytesIO

# 文档解析库
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import markdown
import chardet
from PIL import Image
import pytesseract

from core.ai.config import get_ai_config
logger = logging.getLogger(__name__)

@dataclass
class DocumentChunk:
    """文档块"""
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    chunk_index: int = 0
    start_char: int = 0
    end_char: int = 0
    page_number: Optional[int] = None

@dataclass
class ProcessedDocument:
    """处理后的文档"""
    doc_id: str
    filename: str
    file_type: str
    content: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any] = field(default_factory=dict)
    processed_at: datetime = field(default_factory=datetime.now)
    file_size: int = 0
    page_count: int = 0
    word_count: int = 0
    has_images: bool = False
    images: List[Dict[str, Any]] = field(default_factory=list)

class DocumentProcessor:
    """文档处理器"""

    def __init__(self, config=None):
        self.config = config or get_ai_config()

        # 支持的文件格式和对应的处理器
        self.processors = {
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls,
            '.csv': self._process_csv,
            '.json': self._process_json,
            '.html': self._process_html,
            '.xml': self._process_xml
        }

        # 图片格式
        self.image_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'}

    def process_document(
        self,
        file_path: str,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
        extract_images: Optional[bool] = None
    ) -> ProcessedDocument:
        """处理文档"""
        file_path = Path(file_path)

        # 验证文件
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_size = file_path.stat().st_size
        if file_size > self.config.doc_max_file_size:
            raise ValueError(f"File too large: {file_size} bytes (max: {self.config.doc_max_file_size})")

        # 获取文件类型
        file_ext = file_path.suffix.lower()
        if file_ext not in self.config.doc_supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")

        # 生成文档ID
        doc_id = self._generate_doc_id(file_path)

        # 获取处理器
        processor = self.processors.get(file_ext)
        if not processor:
            raise ValueError(f"No processor for format: {file_ext}")

        # 处理文档
        try:
            logger.info(f"Processing document: {file_path.name}")

            # 调用对应的处理器
            content, metadata = processor(file_path)

            # 处理图片
            images = []
            if (extract_images if extract_images is not None else self.config.doc_extract_images):
                images = self._extract_images_from_content(content, metadata)

            # 分块
            chunk_size = chunk_size or self.config.rag_chunk_size
            chunk_overlap = chunk_overlap or self.config.rag_chunk_overlap
            chunks = self._create_chunks(content, chunk_size, chunk_overlap, metadata)

            # 创建处理后的文档
            processed_doc = ProcessedDocument(
                doc_id=doc_id,
                filename=file_path.name,
                file_type=file_ext,
                content=content,
                chunks=chunks,
                metadata=metadata,
                file_size=file_size,
                page_count=metadata.get('page_count', 1),
                word_count=len(content.split()),
                has_images=len(images) > 0,
                images=images
            )

            logger.info(f"Document processed successfully: {file_path.name}")
            return processed_doc

        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            raise

    def _generate_doc_id(self, file_path: Path) -> str:
        """生成文档ID"""
        # 使用文件路径和修改时间生成唯一ID
        content = f"{file_path.absolute()}:{file_path.stat().st_mtime}"
        return hashlib.sha256(content.encode()).hexdigest()[:16]

    def _process_text(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理文本文件"""
        # 检测编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'

        # 读取内容
        with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()

        metadata = {
            'encoding': encoding,
            'line_count': content.count('\n') + 1
        }

        return content, metadata

    def _process_markdown(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理Markdown文件"""
        content, base_metadata = self._process_text(file_path)

        # 解析Markdown结构
        md = markdown.Markdown(extensions=['meta', 'toc', 'tables'])
        html = md.convert(content)

        # 提取纯文本
        text = re.sub('<[^<]+?>', '', html)

        # 提取标题结构
        headers = re.findall(r'^#{1,6}\s+(.+)$', content, re.MULTILINE)
        metadata = {
            **base_metadata,
            'format': 'markdown',
            'headers': headers,
            'has_code_blocks': '```' in content,
            'has_tables': '|' in content and '---' in content
        }

        return text, metadata

    def _process_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理PDF文件"""
        content_parts = []
        metadata = {
            'format': 'pdf',
            'pages': []
        }

        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            metadata['page_count'] = len(pdf_reader.pages)

            # 提取元数据
            if pdf_reader.metadata:
                metadata['pdf_metadata'] = {
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                }

            # 提取每页内容
            for i, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        content_parts.append(page_text)
                        metadata['pages'].append({
                            'page_number': i + 1,
                            'text_length': len(page_text)
                        })
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {i+1}: {e}")

        content = '\n\n'.join(content_parts)

        # 如果没有提取到文本且启用了OCR
        if not content.strip() and self.config.doc_ocr_enabled:
            logger.info("No text extracted from PDF, trying OCR...")
            content = self._ocr_pdf(file_path)
            metadata['ocr_used'] = True

        return content, metadata

    def _process_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理Word文档"""
        doc = DocxDocument(file_path)
        content_parts = []
        metadata = {
            'format': 'docx',
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables),
            'sections': []
        }

        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)

        # 提取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                if any(row_text):
                    table_text.append(' | '.join(row_text))
            if table_text:
                content_parts.append('\n'.join(table_text))

        # 提取文档属性
        core_props = doc.core_properties
        metadata['properties'] = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'created': core_props.created.isoformat() if core_props.created else None,
            'modified': core_props.modified.isoformat() if core_props.modified else None
        }

        content = '\n\n'.join(content_parts)
        return content, metadata

    def _process_doc(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理旧版Word文档"""
        # 对于.doc格式，可以尝试使用python-docx2txt或调用系统工具
        # 这里简化处理，建议转换为.docx
        logger.warning("Old .doc format detected. Consider converting to .docx for better support.")

        # 尝试作为文本读取（可能会有乱码）
        try:
            return self._process_text(file_path)
        except:
            return "无法解析.doc文件，请转换为.docx格式", {"format": "doc", "error": "unsupported"}

    def _process_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理PowerPoint演示文稿"""
        prs = Presentation(file_path)
        content_parts = []
        metadata = {
            'format': 'pptx',
            'slide_count': len(prs.slides),
            'slides': []
        }

        # 提取每个幻灯片的内容
        for i, slide in enumerate(prs.slides):
            slide_content = []

            # 提取文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            if slide_content:
                slide_text = '\n'.join(slide_content)
                content_parts.append(f"Slide {i+1}:\n{slide_text}")
                metadata['slides'].append({
                    'slide_number': i + 1,
                    'text_length': len(slide_text),
                    'shape_count': len(slide.shapes)
                })

        content = '\n\n'.join(content_parts)
        return content, metadata

    def _process_xlsx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理Excel文件"""
        wb = openpyxl.load_workbook(file_path, read_only=True)
        content_parts = []
        metadata = {
            'format': 'xlsx',
            'sheet_count': len(wb.sheetnames),
            'sheets': []
        }

        # 处理每个工作表
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_content = []

            # 读取数据
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                if any(row_data):
                    sheet_content.append('\t'.join(row_data))

            if sheet_content:
                content_parts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_content))
                metadata['sheets'].append({
                    'name': sheet_name,
                    'row_count': sheet.max_row,
                    'column_count': sheet.max_column
                })

        wb.close()
        content = '\n\n'.join(content_parts)
        return content, metadata

    def _process_xls(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理旧版Excel文件"""
        # 可以使用xlrd库，但新版本不支持.xls
        logger.warning("Old .xls format detected. Consider converting to .xlsx for better support.")
        return "无法解析.xls文件，请转换为.xlsx格式", {"format": "xls", "error": "unsupported"}

    def _process_csv(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理CSV文件"""
        import csv

        # 检测编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'

        content_parts = []
        metadata = {
            'format': 'csv',
            'encoding': encoding,
            'row_count': 0
        }

        with open(file_path, 'r', encoding=encoding) as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                content_parts.append('\t'.join(row))
                metadata['row_count'] += 1

        content = '\n'.join(content_parts)
        return content, metadata

    def _process_json(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理JSON文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # 将JSON转换为可读文本
        content = json.dumps(data, ensure_ascii=False, indent=2)

        metadata = {
            'format': 'json',
            'keys': list(data.keys()) if isinstance(data, dict) else [],
            'type': type(data).__name__
        }

        return content, metadata

    def _process_html(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理HTML文件"""
        from bs4 import BeautifulSoup

        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'html.parser')

        # 提取文本
        text = soup.get_text(separator='\n', strip=True)

        # 提取元数据
        metadata = {
            'format': 'html',
            'title': soup.title.string if soup.title else '',
            'links': len(soup.find_all('a')),
            'images': len(soup.find_all('img')),
            'headers': [h.get_text(strip=True) for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])]
        }

        return text, metadata

    def _process_xml(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """处理XML文件"""
        import xml.etree.ElementTree as ET

        tree = ET.parse(file_path)
        root = tree.getroot()

        # 提取文本内容
        def extract_text(element):
            texts = []
            if element.text and element.text.strip():
                texts.append(element.text.strip())
            for child in element:
                texts.extend(extract_text(child))
            if element.tail and element.tail.strip():
                texts.append(element.tail.strip())
            return texts

        content_parts = extract_text(root)
        content = '\n'.join(content_parts)

        metadata = {
            'format': 'xml',
            'root_tag': root.tag,
            'namespaces': dict(root.attrib.items())
        }

        return content, metadata

    def _create_chunks(
        self,
        content: str,
        chunk_size: int,
        chunk_overlap: int,
        metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """创建文档块"""
        chunks = []

        # 按句子分割，避免在句子中间断开
        sentences = re.split(r'([。！？.!?]\s*)', content)

        current_chunk = []
        current_size = 0
        chunk_index = 0
        start_char = 0

        for i in range(0, len(sentences), 2):
            sentence = sentences[i]
            if i + 1 < len(sentences):
                sentence += sentences[i + 1]
                # 添加标点

            sentence_size = len(sentence)

            # 如果当前块加上新句子超过大小限制
            if current_size + sentence_size > chunk_size and current_chunk:
                # 创建块
                chunk_text = ''.join(current_chunk)
                chunks.append(DocumentChunk(
                    content=chunk_text,
                    metadata={**metadata, 'chunk_method': 'sentence'},
                    chunk_index=chunk_index,
                    start_char=start_char,
                    end_char=start_char + len(chunk_text)
                ))

                # 处理重叠
                if chunk_overlap > 0:
                    # 保留最后的一些内容作为下一块的开始
                    overlap_size = 0
                    overlap_chunks = []
                    for j in range(len(current_chunk) - 1, -1, -1):
                        overlap_size += len(current_chunk[j])
                        overlap_chunks.insert(0, current_chunk[j])
                        if overlap_size >= chunk_overlap:
                            break

                    current_chunk = overlap_chunks
                    current_size = overlap_size
                    start_char = start_char + len(chunk_text) - overlap_size
                else:
                    current_chunk = []
                    current_size = 0
                    start_char = start_char + len(chunk_text)

                chunk_index += 1

            current_chunk.append(sentence)
            current_size += sentence_size

        # 处理最后一块
        if current_chunk:
            chunk_text = ''.join(current_chunk)
            chunks.append(DocumentChunk(
                content=chunk_text,
                metadata={**metadata, 'chunk_method': 'sentence'},
                chunk_index=chunk_index,
                start_char=start_char,
                end_char=start_char + len(chunk_text)
            ))

        return chunks

    def _extract_images_from_content(self, content: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """从内容中提取图片信息"""
        images = []

        # 这里可以根据不同的文档类型实现图片提取
        # 例如从HTML中提取img标签，从Word/PDF中提取嵌入的图片等

        return images

    def _ocr_pdf(self, file_path: Path) -> str:
        """使用OCR处理PDF"""
        try:
            import pdf2image

            # 将PDF转换为图片
            images = pdf2image.convert_from_path(file_path)

            # 对每页进行OCR
            ocr_texts = []
            for i, image in enumerate(images):
                logger.info(f"OCR processing page {i+1}/{len(images)}")
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                ocr_texts.append(text)

            return '\n\n'.join(ocr_texts)

        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""

    def process_image(self, image_path: str) -> ProcessedDocument:
        """处理图片文件（OCR）"""
        image_path = Path(image_path)

        if not self.config.doc_ocr_enabled:
            raise ValueError("OCR is not enabled in configuration")

        # 读取图片
        image = Image.open(image_path)

        # 执行OCR
        text = pytesseract.image_to_string(image, lang='chi_sim+eng')

        # 创建文档
        doc_id = self._generate_doc_id(image_path)

        return ProcessedDocument(
            doc_id=doc_id,
            filename=image_path.name,
            file_type=image_path.suffix,
            content=text,
            chunks=[DocumentChunk(content=text, metadata={'ocr': True})],
            metadata={
                'format': 'image',
                'image_size': image.size,
                'image_mode': image.mode,
                'ocr_used': True
            },
            file_size=image_path.stat().st_size,
            word_count=len(text.split()),
            has_images=True,
            images=[{
                'filename': image_path.name,
                'size': image.size,
                'format': image.format
            }]
        )

    def batch_process(
        self,
        file_paths: List[str],
        **kwargs
    ) -> List[ProcessedDocument]:
        """批量处理文档"""
        results = []

        for file_path in file_paths:
            try:
                doc = self.process_document(file_path, **kwargs)
                results.append(doc)
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")

        return results

    def validate_document(self, doc: ProcessedDocument) -> Tuple[bool, List[str]]:
        """验证处理后的文档"""
        errors = []

        if not doc.content.strip():
            errors.append("文档内容为空")

        if doc.word_count < 10:
            errors.append("文档内容过少")

        if not doc.chunks:
            errors.append("文档未正确分块")

        if doc.file_size == 0:
            errors.append("文件大小为0")

        return len(errors) == 0, errors

# 创建全局实例
_processor_instance: Optional[DocumentProcessor] = None

def get_document_processor() -> DocumentProcessor:
    """获取文档处理器单例"""
    global _processor_instance
    if _processor_instance is None:
        _processor_instance = DocumentProcessor()
    return _processor_instance